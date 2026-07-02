import io
import json
import logging
import os
import shutil
import subprocess
import tempfile
import uuid
import zipfile
from contextlib import asynccontextmanager
from pathlib import Path
from typing import List, Optional, Tuple

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Depends
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.concurrency import run_in_threadpool
from slowapi import _rate_limit_exceeded_handler
from slowapi.errors import RateLimitExceeded

from core.database import connect_db, close_db, get_db
from core.config import settings
from core.deps import get_optional_user
from core.limits import check_file_size, check_tool_access, check_and_increment_ops, get_limits
from core.background import start_background_tasks, cleanup_temp_files
from routers import auth as auth_router
from routers import payments as payments_router
from routers.auth import limiter

logger = logging.getLogger("pdftools")


@asynccontextmanager
async def lifespan(app: FastAPI):
    await connect_db()
    await cleanup_temp_files()  # clean stale files/dirs left over from a crash
    app_url = os.getenv("RENDER_EXTERNAL_URL", "")
    start_background_tasks(app_url)
    yield
    await close_db()


app = FastAPI(title="PDFTools API", version="2.1", lifespan=lifespan)

app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)

app.add_middleware(
    CORSMiddleware,
    allow_origins=[settings.FRONTEND_URL, "http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

app.include_router(auth_router.router)
app.include_router(payments_router.router)

UPLOAD_DIR = Path(tempfile.gettempdir()) / "pdftools"
UPLOAD_DIR.mkdir(exist_ok=True)


# ─── SHARED HELPERS ───────────────────────────────────────────────────────────

def sanitize_filename(name: Optional[str]) -> str:
    """Strip any directory components so a crafted filename (e.g.
    '../../etc/passwd' or an absolute path) can never escape the work dir."""
    base = Path(name or "file").name
    return base or "file"


def new_work_dir() -> Path:
    """A fresh, isolated directory per request. Only used by the LibreOffice
    tools, which need real file paths on disk — everything else stays in
    memory and never touches the filesystem at all."""
    return Path(tempfile.mkdtemp(dir=UPLOAD_DIR, prefix=f"{uuid.uuid4().hex}_"))


async def _read_capped(file: UploadFile, cap: int) -> bytes:
    """Read an upload in chunks, aborting as soon as it exceeds the hard
    ceiling — avoids buffering an arbitrarily large body into memory before
    the plan-based size check even gets a chance to run."""
    chunks = []
    total = 0
    while True:
        chunk = await file.read(1024 * 1024)
        if not chunk:
            break
        total += len(chunk)
        if total > cap:
            raise HTTPException(
                413, f"File exceeds the maximum allowed upload size ({cap // (1024 * 1024)} MB)."
            )
        chunks.append(chunk)
    return b"".join(chunks)


async def _apply_checks(size: int, user, tool_id: str):
    await check_file_size(size, user)
    await check_tool_access(tool_id, user)
    db_ref = get_db()
    if db_ref is not None:
        await check_and_increment_ops(user, db_ref, tool_id)


async def intake_bytes(file: UploadFile, user, tool_id: str) -> bytes:
    """Validate + read a single upload fully into memory. Used by every tool
    that can operate on bytes directly (pypdf/fitz/pdfplumber/etc.) — no
    filesystem interaction, so there's nothing for a crafted filename to do."""
    data = await _read_capped(file, settings.MAX_UPLOAD_BYTES)
    await _apply_checks(len(data), user, tool_id)
    return data


async def intake_bytes_many(files: List[UploadFile], user, tool_id: str) -> List[Tuple[str, bytes]]:
    """Same as intake_bytes() but validates the TOTAL size across all files
    (the original code only checked the first file in a multi-upload)."""
    items = []
    total = 0
    for f in files:
        data = await _read_capped(f, settings.MAX_UPLOAD_BYTES)
        total += len(data)
        if total > settings.MAX_UPLOAD_BYTES:
            raise HTTPException(413, "Combined upload size is too large.")
        items.append((sanitize_filename(f.filename), data))
    await _apply_checks(total, user, tool_id)
    return items


async def intake_to_file(file: UploadFile, user, tool_id: str, work_dir: Path) -> Path:
    """Only for tools whose external process (LibreOffice/wkhtmltopdf) needs
    an actual file path. Filename is sanitized and written inside an
    isolated per-request directory."""
    data = await _read_capped(file, settings.MAX_UPLOAD_BYTES)
    await _apply_checks(len(data), user, tool_id)
    dest = work_dir / sanitize_filename(file.filename)
    dest.write_bytes(data)
    return dest


def pdf_response(data: bytes, filename: str) -> StreamingResponse:
    return StreamingResponse(
        io.BytesIO(data),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


def file_response(data: bytes, media_type: str, filename: str) -> StreamingResponse:
    return StreamingResponse(
        io.BytesIO(data),
        media_type=media_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


def zip_response(data: bytes, filename: str) -> StreamingResponse:
    return StreamingResponse(
        io.BytesIO(data),
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ─── SYNC PROCESSING (run off the event loop via run_in_threadpool) ──────────

def _merge_sync(blobs: List[bytes]) -> bytes:
    from pypdf import PdfWriter, PdfReader
    writer = PdfWriter()
    for data in blobs:
        reader = PdfReader(io.BytesIO(data))
        for page in reader.pages:
            writer.add_page(page)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _parse_ranges(spec: str, total: int) -> set:
    selected = set()
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-")
            selected.update(range(int(a) - 1, min(int(b), total)))
        else:
            selected.add(int(part) - 1)
    return selected


def _split_sync(data: bytes, pages: str) -> bytes:
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(io.BytesIO(data))
    total = len(reader.pages)
    selected = _parse_ranges(pages, total)
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w") as zf:
        for idx in sorted(selected):
            if 0 <= idx < total:
                writer = PdfWriter()
                writer.add_page(reader.pages[idx])
                buf = io.BytesIO()
                writer.write(buf)
                zf.writestr(f"page_{idx + 1}.pdf", buf.getvalue())
    return zip_buf.getvalue()


def _compress_sync(data: bytes) -> bytes:
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    out = doc.tobytes(deflate=True, garbage=4, clean=True)
    doc.close()
    return out


def _rotate_sync(data: bytes, angle: int) -> bytes:
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(io.BytesIO(data))
    writer = PdfWriter()
    for page in reader.pages:
        page.rotate(angle)
        writer.add_page(page)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _remove_pages_sync(data: bytes, pages: str) -> bytes:
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(io.BytesIO(data))
    total = len(reader.pages)
    to_remove = _parse_ranges(pages, total)
    writer = PdfWriter()
    for i, page in enumerate(reader.pages):
        if i not in to_remove:
            writer.add_page(page)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _reorder_sync(data: bytes, order: str) -> bytes:
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(io.BytesIO(data))
    indices = [int(x.strip()) - 1 for x in order.split(",") if x.strip()]
    writer = PdfWriter()
    for idx in indices:
        if 0 <= idx < len(reader.pages):
            writer.add_page(reader.pages[idx])
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _crop_sync(data: bytes, x1: float, y1: float, x2: float, y2: float) -> bytes:
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    out_doc = fitz.open()
    for page in doc:
        # Clamp to THIS page's own rect. Two reasons this matters:
        #  1. Frontend coordinates are Math.round()'d, which can push x2/y2
        #     a fraction of a point past the true page size (e.g. A4's
        #     841.92pt height rounds up to 842) — set_cropbox() rejects
        #     even a 0.01pt overflow.
        #  2. If the doc has mixed page sizes, a clip sized for page 1
        #     could otherwise exceed a smaller page later in the loop.
        pr = page.rect
        cx1 = max(0.0, min(x1, pr.width))
        cy1 = max(0.0, min(y1, pr.height))
        cx2 = max(0.0, min(x2, pr.width))
        cy2 = max(0.0, min(y2, pr.height))
        clip = fitz.Rect(min(cx1, cx2), min(cy1, cy2), max(cx1, cx2), max(cy1, cy2))
        # Guard against a degenerate box (e.g. a page smaller than the
        # requested crop in both dimensions) by falling back to the full page.
        if clip.width < 1 or clip.height < 1:
            clip = pr
        page.set_cropbox(clip)
        out_doc.insert_pdf(doc, from_page=page.number, to_page=page.number)
    result = out_doc.tobytes()
    out_doc.close()
    doc.close()
    return result


def _watermark_sync(data: bytes, text: str, opacity: float) -> bytes:
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    for page in doc:
        w, h = page.rect.width, page.rect.height
        page.insert_text(
            fitz.Point(w * 0.15, h * 0.55), text, fontsize=60,
            color=(0.5, 0.5, 0.5), rotate=45, overlay=True,
        )
    result = doc.tobytes()
    doc.close()
    return result


def _protect_sync(data: bytes, password: str) -> bytes:
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(io.BytesIO(data))
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(password)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _unlock_sync(data: bytes, password: str) -> bytes:
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(io.BytesIO(data))
    if reader.is_encrypted:
        reader.decrypt(password)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _extract_text_sync(data: bytes) -> bytes:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages):
            text_parts.append(f"=== Page {i + 1} ===\n{page.extract_text() or ''}")
    return "\n\n".join(text_parts).encode("utf-8")


def _page_numbers_sync(data: bytes, position: str, start: int) -> bytes:
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    for i, page in enumerate(doc):
        num = start + i
        w, h = page.rect.width, page.rect.height
        if position == "bottom-center":
            pt = fitz.Point(w / 2 - 10, h - 20)
        elif position == "top-center":
            pt = fitz.Point(w / 2 - 10, 30)
        elif position == "bottom-right":
            pt = fitz.Point(w - 50, h - 20)
        else:
            pt = fitz.Point(20, h - 20)
        page.insert_text(pt, str(num), fontsize=11, color=(0, 0, 0))
    result = doc.tobytes()
    doc.close()
    return result


def _pdf_to_jpg_sync(data: bytes, dpi: int) -> bytes:
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w") as zf:
        for i, page in enumerate(doc):
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat)
            zf.writestr(f"page_{i + 1}.jpg", pix.tobytes("jpeg"))
    doc.close()
    return zip_buf.getvalue()


def _jpg_to_pdf_sync(items: List[Tuple[str, bytes]]) -> bytes:
    import fitz
    doc = fitz.open()
    for name, data in items:
        ext = Path(name).suffix.lstrip(".").lower() or "jpg"
        img_doc = fitz.open(stream=data, filetype=ext)
        pdfbytes = img_doc.convert_to_pdf()
        img_doc.close()
        img_pdf = fitz.open("pdf", pdfbytes)
        doc.insert_pdf(img_pdf)
        img_pdf.close()
    result = doc.tobytes()
    doc.close()
    return result


def _pdf_to_word_sync(data: bytes) -> bytes:
    import pdfplumber
    from docx import Document
    doc = Document()
    doc.add_heading("Extracted from PDF", 0)
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages):
            doc.add_heading(f"Page {i + 1}", level=1)
            text = page.extract_text() or ""
            for line in text.split("\n"):
                doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pdf_to_excel_sync(data: bytes) -> bytes:
    import pdfplumber
    import openpyxl
    wb = openpyxl.Workbook()
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages):
            ws = wb.create_sheet(title=f"Page {i + 1}")
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    for row in table:
                        ws.append([cell or "" for cell in row])
            else:
                text = page.extract_text() or ""
                for line in text.split("\n"):
                    ws.append([line])
    if "Sheet" in wb.sheetnames:
        del wb["Sheet"]
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pdf_to_pptx_sync(data: bytes) -> bytes:
    import fitz
    from pptx import Presentation
    from pptx.util import Inches
    doc = fitz.open(stream=data, filetype="pdf")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    for page in doc:
        mat = fitz.Matrix(1.5, 1.5)
        pix = page.get_pixmap(matrix=mat)
        img_buf = io.BytesIO(pix.tobytes("png"))
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_buf, 0, 0, prs.slide_width, prs.slide_height)
    doc.close()
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _libreoffice_to_pdf_sync(work_dir: Path, input_path: Path) -> bytes:
    result = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", str(work_dir), str(input_path)],
        capture_output=True, text=True, timeout=60,
    )
    out = work_dir / (input_path.stem + ".pdf")
    if not out.exists():
        logger.error("LibreOffice conversion failed: %s", result.stderr)
        raise HTTPException(500, "Conversion failed. Please check the uploaded file and try again.")
    return out.read_bytes()


def _html_to_pdf_sync(work_dir: Path, input_path: Path) -> bytes:
    result = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", str(work_dir), str(input_path)],
        capture_output=True, text=True, timeout=60,
    )
    out = work_dir / (input_path.stem + ".pdf")
    if out.exists():
        return out.read_bytes()

    out2 = work_dir / "converted.pdf"
    result2 = subprocess.run(
        ["wkhtmltopdf", str(input_path), str(out2)],
        capture_output=True, text=True, timeout=60,
    )
    if not out2.exists():
        logger.error("HTML to PDF conversion failed: %s / %s", result.stderr, result2.stderr)
        raise HTTPException(500, "HTML to PDF conversion failed.")
    return out2.read_bytes()


def _ocr_sync(data: bytes, page_limit: Optional[int]) -> bytes:
    import fitz
    import pytesseract
    from PIL import Image
    doc = fitz.open(stream=data, filetype="pdf")
    text_parts = []
    for i, page in enumerate(doc):
        if page_limit is not None and i >= page_limit:
            text_parts.append(f"=== Page {i + 1} === [upgrade to Pro to OCR all pages]")
            continue
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        txt = pytesseract.image_to_string(img)
        text_parts.append(f"=== Page {i + 1} ===\n{txt}")
    doc.close()
    return "\n\n".join(text_parts).encode("utf-8")


def _extract_images_sync(data: bytes) -> bytes:
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    zip_buf = io.BytesIO()
    count = 0
    with zipfile.ZipFile(zip_buf, "w") as zf:
        for i, page in enumerate(doc):
            for j, img_info in enumerate(page.get_images(full=True)):
                xref = img_info[0]
                base_image = doc.extract_image(xref)
                zf.writestr(f"page{i + 1}_img{j + 1}.{base_image['ext']}", base_image["image"])
                count += 1
    doc.close()
    if count == 0:
        raise HTTPException(404, "No images found in PDF")
    return zip_buf.getvalue()


def _info_sync(data: bytes) -> dict:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    meta = reader.metadata or {}
    return {
        "pages": len(reader.pages),
        "title": meta.get("/Title", ""),
        "author": meta.get("/Author", ""),
        "subject": meta.get("/Subject", ""),
        "creator": meta.get("/Creator", ""),
        "encrypted": reader.is_encrypted,
    }


# ─── NEW TOOLS: redact, sign, edit, compare, scan-to-pdf, pdf-to-pdfa, repair ─

def _redact_sync(data: bytes, regions: list) -> bytes:
    """regions: list of {page: int (0-indexed), x0, y0, x1, y1} in PDF points.
    Uses PyMuPDF's real redaction (add_redact_annot + apply_redactions), which
    actually strips the underlying text/image content in that area — not just
    a black rectangle drawn on top (which would still leave the original data
    extractable/copyable underneath)."""
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    by_page = {}
    for r in regions:
        by_page.setdefault(int(r["page"]), []).append(r)

    for page_num, rects in by_page.items():
        if page_num < 0 or page_num >= len(doc):
            continue
        page = doc[page_num]
        for r in rects:
            rect = fitz.Rect(r["x0"], r["y0"], r["x1"], r["y1"])
            page.add_redact_annot(rect, fill=(0, 0, 0))
        page.apply_redactions()

    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def _sign_sync(data: bytes, signature_png: bytes, page_num: int, x: float, y: float, width: float, height: float) -> bytes:
    """Stamp a signature image onto a specific page at the given position/size
    (in PDF points, top-left origin, matching how a rendered-page preview in
    the frontend would report click coordinates)."""
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    if page_num < 0 or page_num >= len(doc):
        raise ValueError(f"Page {page_num} out of range (document has {len(doc)} pages)")
    page = doc[page_num]
    rect = fitz.Rect(x, y, x + width, y + height)
    page.insert_image(rect, stream=signature_png)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def _edit_sync(data: bytes, operations: list) -> bytes:
    """operations: list of dicts, each one of:
      {type: "text", page, x, y, text, font_size, color: [r,g,b] 0-1}
      {type: "rect", page, x, y, width, height, color: [r,g,b], fill: bool}
      {type: "circle", page, x, y, radius, color: [r,g,b], fill: bool}
      {type: "line", page, x1, y1, x2, y2, color: [r,g,b]}
    Coordinates in PDF points, top-left origin."""
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")

    for op in operations:
        page_num = int(op["page"])
        if page_num < 0 or page_num >= len(doc):
            continue
        page = doc[page_num]
        color = tuple(op.get("color", [0, 0, 0]))
        op_type = op["type"]

        if op_type == "text":
            page.insert_text(
                (op["x"], op["y"]),
                op.get("text", ""),
                fontsize=float(op.get("font_size", 14)),
                color=color,
            )
        elif op_type == "rect":
            rect = fitz.Rect(op["x"], op["y"], op["x"] + op["width"], op["y"] + op["height"])
            page.draw_rect(rect, color=color, fill=color if op.get("fill") else None)
        elif op_type == "circle":
            center = fitz.Point(op["x"], op["y"])
            page.draw_circle(center, float(op["radius"]), color=color, fill=color if op.get("fill") else None)
        elif op_type == "line":
            page.draw_line(fitz.Point(op["x1"], op["y1"]), fitz.Point(op["x2"], op["y2"]), color=color)

    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def _compare_sync(data_a: bytes, data_b: bytes) -> bytes:
    """Produces a downloadable PDF report: page-by-page unified text diff
    between the two documents. (Text diff, not pixel diff — a 1px rendering
    shift would swamp a visual diff with noise, whereas text diff surfaces
    actual content changes.)"""
    import pdfplumber
    import difflib
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.pagesizes import letter

    def extract_pages(data):
        pages = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for p in pdf.pages:
                pages.append(p.extract_text() or "")
        return pages

    pages_a = extract_pages(data_a)
    pages_b = extract_pages(data_b)
    max_pages = max(len(pages_a), len(pages_b))

    styles = getSampleStyleSheet()
    mono = ParagraphStyle("mono", parent=styles["Normal"], fontName="Courier", fontSize=8, leading=10)
    added_style = ParagraphStyle("added", parent=mono, backColor="#e6ffed")
    removed_style = ParagraphStyle("removed", parent=mono, backColor="#ffeef0")
    heading = ParagraphStyle("heading", parent=styles["Heading2"])

    def esc(s):
        return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    story = [Paragraph("PDF Comparison Report", styles["Title"]), Spacer(1, 12)]
    any_diff = False

    for i in range(max_pages):
        text_a = pages_a[i] if i < len(pages_a) else ""
        text_b = pages_b[i] if i < len(pages_b) else ""
        if text_a == text_b:
            continue
        any_diff = True
        story.append(Paragraph(f"Page {i + 1}", heading))
        lines_a = text_a.splitlines()
        lines_b = text_b.splitlines()
        sm = difflib.SequenceMatcher(None, lines_a, lines_b)
        for tag, i1, i2, j1, j2 in sm.get_opcodes():
            if tag == "equal":
                for line in lines_a[i1:i2][:3]:  # keep unchanged context short
                    story.append(Paragraph(esc(line) or "&nbsp;", mono))
            elif tag == "delete":
                for line in lines_a[i1:i2]:
                    story.append(Paragraph(f"- {esc(line)}", removed_style))
            elif tag == "insert":
                for line in lines_b[j1:j2]:
                    story.append(Paragraph(f"+ {esc(line)}", added_style))
            elif tag == "replace":
                for line in lines_a[i1:i2]:
                    story.append(Paragraph(f"- {esc(line)}", removed_style))
                for line in lines_b[j1:j2]:
                    story.append(Paragraph(f"+ {esc(line)}", added_style))
        story.append(Spacer(1, 16))

    if not any_diff:
        story.append(Paragraph("No text differences were found between the two documents.", styles["Normal"]))
    if len(pages_a) != len(pages_b):
        story.append(Spacer(1, 12))
        story.append(Paragraph(
            f"Note: document A has {len(pages_a)} page(s), document B has {len(pages_b)} page(s).",
            styles["Italic"],
        ))

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter)
    doc.build(story)
    return buf.getvalue()


def _scan_to_pdf_sync(images: List[bytes]) -> bytes:
    """images: list of raw image bytes (phone/desktop camera or gallery
    upload). Applies basic scan cleanup (grayscale + autocontrast) to each,
    then assembles into a single multi-page PDF."""
    from PIL import Image, ImageOps
    pil_pages = []
    for raw in images:
        img = Image.open(io.BytesIO(raw))
        img = img.convert("RGB")
        gray = ImageOps.grayscale(img)
        enhanced = ImageOps.autocontrast(gray, cutoff=1)
        pil_pages.append(enhanced.convert("RGB"))

    if not pil_pages:
        raise ValueError("No images provided")

    buf = io.BytesIO()
    pil_pages[0].save(buf, format="PDF", save_all=True, append_images=pil_pages[1:])
    return buf.getvalue()


def _pdf_to_pdfa_sync(work_dir: Path, src_path: Path) -> bytes:
    """Convert to PDF/A-2b using Ghostscript. Requires the `gs` binary on the
    server — same system-dependency caveat as the LibreOffice-based tools.
    On Render, add `ghostscript` to the apt packages list alongside
    libreoffice/tesseract-ocr/poppler-utils."""
    out_path = work_dir / "converted.pdf"
    cmd = [
        "gs",
        "-dPDFA=2",
        "-dBATCH",
        "-dNOPAUSE",
        "-dNOOUTERSAVE",
        "-sColorConversionStrategy=UseDeviceIndependentColor",
        "-sDEVICE=pdfwrite",
        "-dPDFACompatibilityPolicy=1",
        f"-sOutputFile={out_path}",
        str(src_path),
    ]
    result = subprocess.run(cmd, capture_output=True, timeout=120)
    if result.returncode != 0 or not out_path.exists():
        raise HTTPException(500, f"PDF/A conversion failed: {result.stderr.decode(errors='replace')[:300]}")
    return out_path.read_bytes()


def _repair_sync(data: bytes) -> bytes:
    """Attempt 1: pypdf lenient read + rewrite (fixes many malformed xref
    tables without needing any external process). Attempt 2 (fallback):
    Ghostscript, which recovers more severely damaged PDFs but requires the
    `gs` binary on the server."""
    from pypdf import PdfReader, PdfWriter

    try:
        reader = PdfReader(io.BytesIO(data), strict=False)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        buf = io.BytesIO()
        writer.write(buf)
        result = buf.getvalue()
        # Sanity check: does the rewritten file actually open and have pages?
        PdfReader(io.BytesIO(result), strict=False).pages[0]
        return result
    except Exception:
        pass  # fall through to Ghostscript

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        src = tmp_path / "input.pdf"
        src.write_bytes(data)
        out = tmp_path / "repaired.pdf"
        cmd = ["gs", "-o", str(out), "-sDEVICE=pdfwrite", "-dPDFSTOPONERROR=false", str(src)]
        result = subprocess.run(cmd, capture_output=True, timeout=120)
        if result.returncode != 0 or not out.exists():
            raise HTTPException(500, f"Could not repair this PDF: {result.stderr.decode(errors='replace')[:300]}")
        return out.read_bytes()



# ─── ROUTES ───────────────────────────────────────────────────────────────────

@app.post("/merge")
async def merge_pdfs(files: List[UploadFile] = File(...), user=Depends(get_optional_user)):
    items = await intake_bytes_many(files, user, "merge")
    result = await run_in_threadpool(_merge_sync, [data for _, data in items])
    return pdf_response(result, "merged.pdf")


@app.post("/split")
async def split_pdf(file: UploadFile = File(...), pages: str = Form(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "split")
    result = await run_in_threadpool(_split_sync, data, pages)
    return zip_response(result, "split_pages.zip")


@app.post("/compress")
async def compress_pdf(file: UploadFile = File(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "compress")
    result = await run_in_threadpool(_compress_sync, data)
    return pdf_response(result, "compressed.pdf")


@app.post("/rotate")
async def rotate_pdf(file: UploadFile = File(...), angle: int = Form(90), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "rotate")
    result = await run_in_threadpool(_rotate_sync, data, angle)
    return pdf_response(result, "rotated.pdf")


@app.post("/remove-pages")
async def remove_pages(file: UploadFile = File(...), pages: str = Form(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "remove-pages")
    result = await run_in_threadpool(_remove_pages_sync, data, pages)
    return pdf_response(result, "removed_pages.pdf")


@app.post("/reorder")
async def reorder_pages(file: UploadFile = File(...), order: str = Form(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "reorder")
    result = await run_in_threadpool(_reorder_sync, data, order)
    return pdf_response(result, "reordered.pdf")


@app.post("/crop")
async def crop_pdf(
    file: UploadFile = File(...),
    x1: float = Form(0), y1: float = Form(0),
    x2: float = Form(595), y2: float = Form(842),
    user=Depends(get_optional_user),
):
    data = await intake_bytes(file, user, "crop")
    result = await run_in_threadpool(_crop_sync, data, x1, y1, x2, y2)
    return pdf_response(result, "cropped.pdf")


@app.post("/watermark")
async def watermark_pdf(
    file: UploadFile = File(...),
    text: str = Form("WATERMARK"),
    opacity: float = Form(0.3),
    user=Depends(get_optional_user),
):
    data = await intake_bytes(file, user, "watermark")
    result = await run_in_threadpool(_watermark_sync, data, text, opacity)
    return pdf_response(result, "watermarked.pdf")


@app.post("/protect")
async def protect_pdf(file: UploadFile = File(...), password: str = Form(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "protect")
    result = await run_in_threadpool(_protect_sync, data, password)
    return pdf_response(result, "protected.pdf")


@app.post("/unlock")
async def unlock_pdf(file: UploadFile = File(...), password: str = Form(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "unlock")
    result = await run_in_threadpool(_unlock_sync, data, password)
    return pdf_response(result, "unlocked.pdf")


@app.post("/extract-text")
async def extract_text(file: UploadFile = File(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "extract-text")
    result = await run_in_threadpool(_extract_text_sync, data)
    return file_response(result, "text/plain", "extracted_text.txt")


@app.post("/page-numbers")
async def add_page_numbers(
    file: UploadFile = File(...),
    position: str = Form("bottom-center"),
    start: int = Form(1),
    user=Depends(get_optional_user),
):
    data = await intake_bytes(file, user, "page-numbers")
    result = await run_in_threadpool(_page_numbers_sync, data, position, start)
    return pdf_response(result, "numbered.pdf")


@app.post("/pdf-to-jpg")
async def pdf_to_jpg(file: UploadFile = File(...), dpi: int = Form(150), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "pdf-to-jpg")
    result = await run_in_threadpool(_pdf_to_jpg_sync, data, dpi)
    return zip_response(result, "pdf_pages.zip")


@app.post("/jpg-to-pdf")
async def jpg_to_pdf(files: List[UploadFile] = File(...), user=Depends(get_optional_user)):
    items = await intake_bytes_many(files, user, "jpg-to-pdf")
    result = await run_in_threadpool(_jpg_to_pdf_sync, items)
    return pdf_response(result, "images.pdf")


@app.post("/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "pdf-to-word")
    result = await run_in_threadpool(_pdf_to_word_sync, data)
    return file_response(
        result,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "converted.docx",
    )


@app.post("/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "pdf-to-excel")
    result = await run_in_threadpool(_pdf_to_excel_sync, data)
    return file_response(
        result,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "converted.xlsx",
    )


@app.post("/pdf-to-pptx")
async def pdf_to_pptx(file: UploadFile = File(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "pdf-to-pptx")
    result = await run_in_threadpool(_pdf_to_pptx_sync, data)
    return file_response(
        result,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "converted.pptx",
    )


@app.post("/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...), user=Depends(get_optional_user)):
    work_dir = new_work_dir()
    try:
        path = await intake_to_file(file, user, "word-to-pdf", work_dir)
        result = await run_in_threadpool(_libreoffice_to_pdf_sync, work_dir, path)
        return pdf_response(result, "converted.pdf")
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)


@app.post("/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...), user=Depends(get_optional_user)):
    work_dir = new_work_dir()
    try:
        path = await intake_to_file(file, user, "excel-to-pdf", work_dir)
        result = await run_in_threadpool(_libreoffice_to_pdf_sync, work_dir, path)
        return pdf_response(result, "converted.pdf")
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)


@app.post("/pptx-to-pdf")
async def pptx_to_pdf(file: UploadFile = File(...), user=Depends(get_optional_user)):
    work_dir = new_work_dir()
    try:
        path = await intake_to_file(file, user, "pptx-to-pdf", work_dir)
        result = await run_in_threadpool(_libreoffice_to_pdf_sync, work_dir, path)
        return pdf_response(result, "converted.pdf")
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)


@app.post("/html-to-pdf")
async def html_to_pdf(file: UploadFile = File(...), user=Depends(get_optional_user)):
    work_dir = new_work_dir()
    try:
        path = await intake_to_file(file, user, "html-to-pdf", work_dir)
        result = await run_in_threadpool(_html_to_pdf_sync, work_dir, path)
        return pdf_response(result, "converted.pdf")
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)


@app.post("/ocr")
async def ocr_pdf(file: UploadFile = File(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "ocr")
    limits = get_limits(user)
    result = await run_in_threadpool(_ocr_sync, data, limits["ocr_page_limit"])
    return file_response(result, "text/plain", "ocr_output.txt")


@app.post("/extract-images")
async def extract_images(file: UploadFile = File(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "extract-images")
    result = await run_in_threadpool(_extract_images_sync, data)
    return zip_response(result, "extracted_images.zip")


@app.post("/info")
async def pdf_info(file: UploadFile = File(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "info")
    return await run_in_threadpool(_info_sync, data)


@app.post("/redact")
async def redact_pdf(file: UploadFile = File(...), regions: str = Form(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "redact")
    try:
        parsed_regions = json.loads(regions)
    except (TypeError, ValueError):
        raise HTTPException(400, 'regions must be valid JSON, e.g. [{"page":0,"x0":10,"y0":10,"x1":100,"y1":40}]')
    result = await run_in_threadpool(_redact_sync, data, parsed_regions)
    return pdf_response(result, "redacted.pdf")


@app.post("/edit")
async def edit_pdf(file: UploadFile = File(...), operations: str = Form(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "edit")
    try:
        parsed_ops = json.loads(operations)
    except (TypeError, ValueError):
        raise HTTPException(400, 'operations must be valid JSON, e.g. [{"type":"text","page":0,"x":100,"y":400,"text":"Hi"}]')
    result = await run_in_threadpool(_edit_sync, data, parsed_ops)
    return pdf_response(result, "edited.pdf")


@app.post("/sign")
async def sign_pdf(
    file: UploadFile = File(...),
    signature: UploadFile = File(...),
    page: int = Form(0),
    x: float = Form(...),
    y: float = Form(...),
    width: float = Form(150),
    height: float = Form(60),
    user=Depends(get_optional_user),
):
    data = await intake_bytes(file, user, "sign")
    sig_bytes = await signature.read()
    if len(sig_bytes) > settings.MAX_UPLOAD_BYTES:
        raise HTTPException(413, "Signature image is too large.")
    result = await run_in_threadpool(_sign_sync, data, sig_bytes, page, x, y, width, height)
    return pdf_response(result, "signed.pdf")


@app.post("/compare")
async def compare_pdfs(files: List[UploadFile] = File(...), user=Depends(get_optional_user)):
    if len(files) != 2:
        raise HTTPException(400, "Upload exactly two PDFs to compare (original first, then the revised version).")
    items = await intake_bytes_many(files, user, "compare")
    result = await run_in_threadpool(_compare_sync, items[0][1], items[1][1])
    return pdf_response(result, "comparison_report.pdf")


@app.post("/scan-to-pdf")
async def scan_to_pdf(files: List[UploadFile] = File(...), user=Depends(get_optional_user)):
    items = await intake_bytes_many(files, user, "scan-to-pdf")
    result = await run_in_threadpool(_scan_to_pdf_sync, [data for _, data in items])
    return pdf_response(result, "scanned.pdf")


@app.post("/pdf-to-pdfa")
async def pdf_to_pdfa(file: UploadFile = File(...), user=Depends(get_optional_user)):
    work_dir = new_work_dir()
    try:
        path = await intake_to_file(file, user, "pdf-to-pdfa", work_dir)
        result = await run_in_threadpool(_pdf_to_pdfa_sync, work_dir, path)
        return pdf_response(result, "converted_pdfa.pdf")
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)


@app.post("/repair")
async def repair_pdf(file: UploadFile = File(...), user=Depends(get_optional_user)):
    data = await intake_bytes(file, user, "repair")
    result = await run_in_threadpool(_repair_sync, data)
    return pdf_response(result, "repaired.pdf")


@app.get("/health")
def health():
    return {"status": "ok"}
